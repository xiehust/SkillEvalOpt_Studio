"""Bounded inspection and rendering for Word and PowerPoint artifacts."""
from __future__ import annotations

import datetime as dt
import json
import math
import os
import re
import shutil
import stat
import struct
import tempfile
import time
import unicodedata
import zipfile
import zlib
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import unquote_to_bytes, urlsplit
from xml.etree import ElementTree

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from skillopt.envs.skilleval import artifacts as artifact_security

from ._scratch import current_scratch_transaction
from .base import (
    InspectionError,
    RenderBudget,
    ResponseBudget,
    bounded_diagnostic,
    safe_run,
    validate_json_result,
)
from . import spreadsheet as ooxml

_OLE_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_WORD_MAIN_PART = "word/document.xml"
_PPT_MAIN_PART = "ppt/presentation.xml"
_WORD_MAIN_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "wordprocessingml.document.main+xml"
)
_PPT_MAIN_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "presentationml.presentation.main+xml"
)
_WORD_HEADER_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "wordprocessingml.header+xml"
)
_WORD_FOOTER_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "wordprocessingml.footer+xml"
)
_OFFICE_DOCUMENT_RELATIONSHIP = (
    "http://schemas.openxmlformats.org/officeDocument/2006/"
    "relationships/officeDocument"
)
_OFFICE_RELATIONSHIPS_NAMESPACE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)
_PACKAGE_RELATIONSHIPS_NAMESPACE = (
    "http://schemas.openxmlformats.org/package/2006/relationships"
)
_WORD_NAMESPACE = (
    "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
)
_PRESENTATION_NAMESPACE = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)
_DRAWING_NAMESPACE = (
    "http://schemas.openxmlformats.org/drawingml/2006/main"
)
_CHART_NAMESPACE = (
    "http://schemas.openxmlformats.org/drawingml/2006/chart"
)
_CHART_DRAWING_NAMESPACE = (
    "http://schemas.openxmlformats.org/drawingml/2006/chartDrawing"
)
_DIAGRAM_NAMESPACE = (
    "http://schemas.openxmlformats.org/drawingml/2006/diagram"
)
_DIAGRAM_DRAWING_NAMESPACE = (
    "http://schemas.microsoft.com/office/drawing/2008/diagram"
)
_SPREADSHEET_DRAWING_NAMESPACE = (
    "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
)
_CHART_STYLE_NAMESPACE = (
    "http://schemas.microsoft.com/office/drawing/2012/chartStyle"
)
_CHARTEX_NAMESPACE = (
    "http://schemas.microsoft.com/office/drawing/2014/chartex"
)
_CORE_PROPERTIES_NAMESPACE = (
    "http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
)
_EXTENDED_PROPERTIES_NAMESPACE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/"
    "extended-properties"
)
_CUSTOM_PROPERTIES_NAMESPACE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/"
    "custom-properties"
)
_CUSTOM_XML_NAMESPACE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/customXml"
)
_XML_SIGNATURE_NAMESPACE = "http://www.w3.org/2000/09/xmldsig#"
_INK_NAMESPACE = "http://www.w3.org/2003/InkML"
_HEADER_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/header"
_FOOTER_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/footer"
_SLIDE_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/slide"
_SLIDE_LAYOUT_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/slideLayout"
)
_NOTES_SLIDE_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/notesSlide"
)
_IMAGE_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/image"
_AUDIO_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/audio"
_VIDEO_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/video"
_CHART_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/chart"
_MEDIA_RELATIONSHIP = (
    "http://schemas.microsoft.com/office/2007/relationships/media"
)
_SLIDE_MASTER_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/slideMaster"
)
_NOTES_MASTER_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/notesMaster"
)
_HANDOUT_MASTER_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/handoutMaster"
)
_PRES_PROPS_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/presProps"
_VIEW_PROPS_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/viewProps"
_TABLE_STYLES_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/tableStyles"
)
_THEME_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/theme"
_THEME_OVERRIDE_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/themeOverride"
)
_COMMENTS_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/comments"
_COMMENT_AUTHORS_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/commentAuthors"
)
_TAGS_RELATIONSHIP = f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/tags"
_CHART_USER_SHAPES_RELATIONSHIP = (
    f"{_OFFICE_RELATIONSHIPS_NAMESPACE}/chartUserShapes"
)
_SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "presentationml.slide+xml"
)
_SLIDE_LAYOUT_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "presentationml.slideLayout+xml"
)
_NOTES_SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "presentationml.notesSlide+xml"
)
_CHART_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
    "application/vnd.ms-office.chartex+xml",
}
_PPT_RELATIONSHIP_TARGET_ROOTS = {
    "application/vnd.openxmlformats-package.core-properties+xml": (
        f"{{{_CORE_PROPERTIES_NAMESPACE}}}coreProperties"
    ),
    "application/vnd.openxmlformats-officedocument."
    "extended-properties+xml": (
        f"{{{_EXTENDED_PROPERTIES_NAMESPACE}}}Properties"
    ),
    "application/vnd.openxmlformats-officedocument."
    "custom-properties+xml": (
        f"{{{_CUSTOM_PROPERTIES_NAMESPACE}}}Properties"
    ),
    "application/vnd.openxmlformats-officedocument."
    "customXmlProperties+xml": (
        f"{{{_CUSTOM_XML_NAMESPACE}}}datastoreItem"
    ),
    _PPT_MAIN_CONTENT_TYPE: (
        f"{{{_PRESENTATION_NAMESPACE}}}presentation"
    ),
    _SLIDE_CONTENT_TYPE: f"{{{_PRESENTATION_NAMESPACE}}}sld",
    _SLIDE_LAYOUT_CONTENT_TYPE: (
        f"{{{_PRESENTATION_NAMESPACE}}}sldLayout"
    ),
    _NOTES_SLIDE_CONTENT_TYPE: (
        f"{{{_PRESENTATION_NAMESPACE}}}notes"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.slideMaster+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}sldMaster"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.notesMaster+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}notesMaster"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.handoutMaster+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}handoutMaster"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.presProps+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}presentationPr"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.viewProps+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}viewPr"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.tableStyles+xml": (
        f"{{{_DRAWING_NAMESPACE}}}tblStyleLst"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.comments+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}cmLst"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.commentAuthors+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}cmAuthorLst"
    ),
    "application/vnd.openxmlformats-officedocument."
    "presentationml.tags+xml": (
        f"{{{_PRESENTATION_NAMESPACE}}}tagLst"
    ),
    "application/vnd.openxmlformats-officedocument.theme+xml": (
        f"{{{_DRAWING_NAMESPACE}}}theme"
    ),
    "application/vnd.openxmlformats-officedocument.themeOverride+xml": (
        f"{{{_DRAWING_NAMESPACE}}}themeOverride"
    ),
    "application/vnd.openxmlformats-officedocument."
    "drawingml.chart+xml": f"{{{_CHART_NAMESPACE}}}chartSpace",
    "application/vnd.ms-office.chartex+xml": (
        f"{{{_CHARTEX_NAMESPACE}}}chartSpace"
    ),
    "application/vnd.openxmlformats-officedocument."
    "drawingml.chartshapes+xml": (
        f"{{{_CHART_DRAWING_NAMESPACE}}}userShapes"
    ),
    "application/vnd.openxmlformats-officedocument."
    "drawingml.diagramColors+xml": (
        f"{{{_DIAGRAM_NAMESPACE}}}colorsDef"
    ),
    "application/vnd.openxmlformats-officedocument."
    "drawingml.diagramData+xml": (
        f"{{{_DIAGRAM_NAMESPACE}}}dataModel"
    ),
    "application/vnd.openxmlformats-officedocument."
    "drawingml.diagramLayout+xml": (
        f"{{{_DIAGRAM_NAMESPACE}}}layoutDef"
    ),
    "application/vnd.openxmlformats-officedocument."
    "drawingml.diagramStyle+xml": (
        f"{{{_DIAGRAM_NAMESPACE}}}styleDef"
    ),
    "application/vnd.ms-office.drawingml.diagramDrawing+xml": (
        f"{{{_DIAGRAM_DRAWING_NAMESPACE}}}drawing"
    ),
    "application/vnd.openxmlformats-officedocument.drawing+xml": (
        f"{{{_SPREADSHEET_DRAWING_NAMESPACE}}}wsDr"
    ),
    "application/vnd.ms-office.chartstyle+xml": (
        f"{{{_CHART_STYLE_NAMESPACE}}}chartStyle"
    ),
    "application/vnd.ms-office.chartcolorstyle+xml": (
        f"{{{_CHART_STYLE_NAMESPACE}}}colorStyle"
    ),
    "application/vnd.openxmlformats-package."
    "digital-signature-xmlsignature+xml": (
        f"{{{_XML_SIGNATURE_NAMESPACE}}}Signature"
    ),
    "application/inkml+xml": f"{{{_INK_NAMESPACE}}}ink",
}
_WORD_RELATIONSHIP_TARGET_ROOTS = {
    _HEADER_RELATIONSHIP: (
        _WORD_HEADER_CONTENT_TYPE,
        f"{{{_WORD_NAMESPACE}}}hdr",
    ),
    _FOOTER_RELATIONSHIP: (
        _WORD_FOOTER_CONTENT_TYPE,
        f"{{{_WORD_NAMESPACE}}}ftr",
    ),
}
_PPT_RELATIONSHIP_TYPE_CONTENT_TYPES = {
    _SLIDE_MASTER_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.slideMaster+xml"
    ),
    _NOTES_MASTER_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.notesMaster+xml"
    ),
    _HANDOUT_MASTER_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.handoutMaster+xml"
    ),
    _PRES_PROPS_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.presProps+xml"
    ),
    _VIEW_PROPS_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.viewProps+xml"
    ),
    _TABLE_STYLES_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.tableStyles+xml"
    ),
    _THEME_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument.theme+xml"
    ),
    _THEME_OVERRIDE_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument.themeOverride+xml"
    ),
    _COMMENTS_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.comments+xml"
    ),
    _COMMENT_AUTHORS_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.commentAuthors+xml"
    ),
    _TAGS_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.tags+xml"
    ),
    _CHART_USER_SHAPES_RELATIONSHIP: (
        "application/vnd.openxmlformats-officedocument."
        "drawingml.chartshapes+xml"
    ),
}
_RELATIONSHIPS_CONTENT_TYPE = (
    "application/vnd.openxmlformats-package.relationships+xml"
)
_MAX_XML_NODES = 1_000_000
_MAX_XML_DEPTH = 256
_MAX_PARAGRAPHS = 250_000
_MAX_RUNS = 1_000_000
_MAX_TABLE_CELLS = 250_000
_MAX_SLIDES = 10_000
_MAX_SHAPES = 250_000
_MAX_NOTES = 10_000
_MAX_RELATIONSHIPS = 100_000
_MAX_MEDIA_PART_BYTES = 64 * 1024 * 1024
_MAX_INDEX_ITEMS = 128
_MAX_METADATA_ITEMS = 128
_MAX_TEXT_CHARS = 4_096
_MAX_SHAPE_TEXT_ITEMS = 16
_PARAGRAPH_PAGE_SIZE = 64
_TABLE_PAGE_SIZE = 16
_HEADER_PAGE_SIZE = 16
_SHAPE_PAGE_SIZE = 64
_TABLE_CELL_PAGE_SIZE = 64
_TABLE_ROW_PAGE_SIZE = 16
_LIBREOFFICE_TIMEOUT_SECONDS = 180
_LIBREOFFICE_MAX_TRANSIENT_ATTEMPTS = 3
_ZIP_STREAM_CHUNK = 1024 * 1024
_PAGE_SELECTOR_RE = re.compile(r"^page:[1-9][0-9]*$")
_DOC_PAGE_SELECTOR_RE = re.compile(
    r"^(paragraphs|tables|headers|footers):page:([1-9][0-9]*)$"
)
_DOC_ITEM_SELECTOR_RE = re.compile(
    r"^(paragraph|table|header|footer|section):([1-9][0-9]*)"
    r"(?::page:([1-9][0-9]*))?$"
)
_DOC_SECTIONS_CURSOR_RE = re.compile(
    r"^sections:cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_DOC_HEADER_TABLES_CURSOR_RE = re.compile(
    r"^(header|footer):([1-9][0-9]*):tables:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_DOC_TABLE_ROWS_CURSOR_RE = re.compile(
    r"^table:([1-9][0-9]*):rows:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_DOC_TABLE_CELLS_CURSOR_RE = re.compile(
    r"^table:([1-9][0-9]*):row:([1-9][0-9]*):cells:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_DOC_SCOPED_TABLE_ROWS_CURSOR_RE = re.compile(
    r"^(header|footer):([1-9][0-9]*):table:([1-9][0-9]*):"
    r"rows:cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_DOC_SCOPED_TABLE_CELLS_CURSOR_RE = re.compile(
    r"^(header|footer):([1-9][0-9]*):table:([1-9][0-9]*):"
    r"row:([1-9][0-9]*):cells:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_SLIDE_SELECTOR_RE = re.compile(
    r"^slide:([1-9][0-9]*)(?::page:([1-9][0-9]*))?$"
)
_SLIDE_SHAPES_CURSOR_RE = re.compile(
    r"^slide:([1-9][0-9]*):shapes:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_SLIDE_CHILDREN_CURSOR_RE = re.compile(
    r"^slide:([1-9][0-9]*):shape:"
    r"([1-9][0-9]*(?:\.[1-9][0-9]*)*):children:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_SLIDE_NOTES_CURSOR_RE = re.compile(
    r"^slide:([1-9][0-9]*):notes:"
    r"cursor:(0|[1-9][0-9]*):limit:([1-9][0-9]*)$"
)
_OLE_FREESECT = 0xFFFFFFFF
_OLE_ENDOFCHAIN = 0xFFFFFFFE
_OLE_FATSECT = 0xFFFFFFFD
_OLE_DIFSECT = 0xFFFFFFFC


@dataclass(frozen=True)
class _PackageInfo:
    kind: str
    relationships: tuple[dict, ...]
    media: tuple[dict, ...]
    embedded_objects: tuple[dict, ...]
    macros: tuple[dict, ...]


def _json_bytes(value: object) -> int:
    return len(
        json.dumps(
            value,
            ensure_ascii=False,
            allow_nan=False,
            separators=(",", ":"),
        ).encode("utf-8")
    )


def _text_chars(value: object) -> int:
    if isinstance(value, str):
        return len(value)
    if isinstance(value, list):
        return sum(_text_chars(item) for item in value)
    if isinstance(value, dict):
        return sum(
            len(key) + _text_chars(child)
            for key, child in value.items()
        )
    return 0


def _bounded_text(value: object, limit: int = _MAX_TEXT_CHARS) -> str:
    text = "" if value is None else str(value)
    return text if len(text) <= limit else text[:limit]


def _json_scalar(value: object):
    if value is None or isinstance(value, (str, bool, int)):
        return value
    if isinstance(value, float):
        return value if math.isfinite(value) else str(value)
    if isinstance(value, (dt.datetime, dt.date, dt.time)):
        return value.isoformat()
    return _bounded_text(value)


def _metadata_page(items: list[dict], budget: ResponseBudget) -> dict:
    page = {"items": [], "returned": 0, "omitted": len(items)}
    for item in items[:_MAX_METADATA_ITEMS]:
        candidate = {
            "items": [*page["items"], item],
            "returned": page["returned"] + 1,
            "omitted": len(items) - page["returned"] - 1,
        }
        if _json_bytes(candidate) > max(0, budget.max_bytes // 3):
            break
        page = candidate
    return page


class _OfficeXmlValidator:
    def __init__(
        self,
        name: str,
        content_type: str | None,
        totals: dict[str, int],
    ) -> None:
        self._name = name
        self._content_type = (content_type or "").casefold()
        self._totals = totals
        self._parser = ElementTree.XMLPullParser(events=("start", "end"))
        self._depth = 0
        self._bytes = 0
        self._tail = b""

    def _increment(self, key: str, limit: int, label: str) -> None:
        value = self._totals.get(key, 0) + 1
        self._totals[key] = value
        if value > limit:
            raise InspectionError(f"OOXML {label} count exceeds configured limit")

    def _drain(self) -> None:
        for event, element in self._parser.read_events():
            if event == "start":
                self._depth += 1
                if self._depth > _MAX_XML_DEPTH:
                    raise InspectionError("OOXML XML depth exceeds configured limit")
                self._increment("nodes", _MAX_XML_NODES, "XML node")
                local = (
                    element.tag.rsplit("}", 1)[-1]
                    if isinstance(element.tag, str)
                    else ""
                )
                if "wordprocessingml." in self._content_type:
                    if local == "p":
                        self._increment(
                            "paragraphs", _MAX_PARAGRAPHS, "paragraph"
                        )
                    elif local == "r":
                        self._increment("runs", _MAX_RUNS, "run")
                    elif local == "tc":
                        self._increment(
                            "table_cells", _MAX_TABLE_CELLS, "table cell"
                        )
                if self._content_type == _PPT_MAIN_CONTENT_TYPE and local == "sldId":
                    self._increment("slides", _MAX_SLIDES, "slide")
                if ".presentationml.slide+xml" in self._content_type and local in {
                    "sp",
                    "grpSp",
                    "graphicFrame",
                    "pic",
                    "cxnSp",
                }:
                    self._increment("shapes", _MAX_SHAPES, "shape")
                if ".presentationml.notesSlide+xml" in self._content_type:
                    if local == "cSld":
                        self._increment("notes", _MAX_NOTES, "notes slide")
            else:
                self._depth -= 1
                if self._depth < 0:
                    raise InspectionError("OOXML XML part has invalid nesting")
                element.clear()

    def feed(self, chunk: bytes) -> None:
        self._bytes += len(chunk)
        if self._bytes > ooxml._MAX_OOXML_XML_PART_BYTES:
            raise InspectionError(
                f"OOXML XML part bytes exceed configured limit: "
                f"{bounded_diagnostic(self._name)!r}"
            )
        declaration_scan = (self._tail + chunk).lower()
        if b"<!doctype" in declaration_scan or b"<!entity" in declaration_scan:
            raise InspectionError(
                "OOXML XML part contains forbidden declarations"
            )
        self._tail = declaration_scan[-16:]
        try:
            self._parser.feed(chunk)
            self._drain()
        except ElementTree.ParseError as exc:
            raise InspectionError(
                f"OOXML XML part is malformed: "
                f"{bounded_diagnostic(self._name)!r}"
            ) from exc

    def finish(self) -> None:
        try:
            self._parser.close()
            self._drain()
        except ElementTree.ParseError as exc:
            raise InspectionError(
                f"OOXML XML part is malformed: "
                f"{bounded_diagnostic(self._name)!r}"
            ) from exc
        if self._depth != 0:
            raise InspectionError("OOXML XML part has invalid nesting")


def _relationship_source(name: str) -> str:
    if name == "_rels/.rels":
        return ""
    parent, marker, filename = name.rpartition("/_rels/")
    if not marker or not filename.endswith(".rels"):
        raise InspectionError("OOXML relationship part path is invalid")
    source = f"{parent}/{filename[:-5]}" if parent else filename[:-5]
    if not artifact_security._safe_zip_member(source):
        raise InspectionError("OOXML relationship source is invalid")
    return source


def _relationship_target(source: str, target: str) -> str:
    if not target or "\x00" in target or "\\" in target:
        raise InspectionError("OOXML relationship target is invalid")
    split = urlsplit(target)
    if split.scheme or split.netloc or split.query or split.fragment:
        raise InspectionError("OOXML relationship target is not an internal part")
    try:
        decoded = unquote_to_bytes(split.path).decode("utf-8")
    except UnicodeDecodeError as exc:
        raise InspectionError(
            "OOXML relationship target encoding is invalid"
        ) from exc
    parts = [] if decoded.startswith("/") else source.split("/")[:-1]
    for part in decoded.lstrip("/").split("/"):
        if part in {"", "."}:
            continue
        if part == "..":
            if not parts:
                raise InspectionError(
                    "OOXML relationship target escapes the package"
                )
            parts.pop()
        else:
            parts.append(part)
    normalized = "/".join(parts)
    if not artifact_security._safe_zip_member(normalized):
        raise InspectionError("OOXML relationship target is invalid")
    return normalized


def _parse_relationship_part(
    archive: zipfile.ZipFile,
    member: zipfile.ZipInfo,
    source: str,
    names: set[str],
    totals: dict[str, int],
) -> list[dict]:
    payload = ooxml._read_bounded_member(
        archive,
        member,
        ooxml._MAX_OOXML_GRAPH_PART_BYTES,
        "OOXML relationship part",
    )
    lowered = payload.lower()
    if b"<!doctype" in lowered or b"<!entity" in lowered:
        raise InspectionError(
            "OOXML relationship part contains forbidden declarations"
        )
    try:
        root = ElementTree.fromstring(payload)
    except ElementTree.ParseError as exc:
        raise InspectionError("OOXML relationship part is malformed") from exc
    relationships_tag = (
        f"{{{_PACKAGE_RELATIONSHIPS_NAMESPACE}}}Relationships"
    )
    relationship_tag = (
        f"{{{_PACKAGE_RELATIONSHIPS_NAMESPACE}}}Relationship"
    )
    if root.tag != relationships_tag or root.attrib:
        raise InspectionError("OOXML relationship root is invalid")
    rows: list[dict] = []
    ids: set[str] = set()
    for element in root:
        if (
            element.tag != relationship_tag
            or len(element)
            or not set(element.attrib).issubset(
                {"Id", "Type", "Target", "TargetMode"}
            )
        ):
            raise InspectionError("OOXML relationship entry is invalid")
        relationship_id = element.attrib.get("Id", "")
        relationship_type = element.attrib.get("Type", "")
        target = element.attrib.get("Target", "")
        target_mode = element.attrib.get("TargetMode")
        if (
            not relationship_id
            or relationship_id in ids
            or not relationship_type
            or not target
            or target_mode not in {None, "External"}
            or max(map(len, (relationship_id, relationship_type, target)))
            > ooxml._MAX_OOXML_GRAPH_PART_BYTES
        ):
            raise InspectionError("OOXML relationship entry is invalid")
        ids.add(relationship_id)
        totals["relationships"] = totals.get("relationships", 0) + 1
        if totals["relationships"] > _MAX_RELATIONSHIPS:
            raise InspectionError(
                "OOXML relationship count exceeds configured limit"
            )
        external = target_mode == "External"
        normalized = None if external else _relationship_target(source, target)
        if normalized is not None and normalized not in names:
            raise InspectionError("OOXML relationship target part is missing")
        rows.append(
            {
                "source": source or "/",
                "id": _bounded_text(relationship_id, 256),
                "type": _bounded_text(relationship_type, 512),
                "target": _bounded_text(target, 1_024),
                "external": external,
                "_normalized_target": normalized,
            }
        )
    return rows


def _validate_ppt_relationship_target_roots(
    archive: zipfile.ZipFile,
    members: dict[str, zipfile.ZipInfo],
    part_content_types: dict[str, str],
    relationships: list[dict],
) -> dict[str, str]:
    root_tags: dict[str, str] = {}
    for relationship in relationships:
        if relationship["external"]:
            continue
        target = relationship.get("_normalized_target")
        if not isinstance(target, str):
            raise InspectionError(
                "OOXML PowerPoint relationship target is invalid"
            )
        content_type = part_content_types.get(target)
        expected_content_type = _PPT_RELATIONSHIP_TYPE_CONTENT_TYPES.get(
            relationship["type"]
        )
        if (
            expected_content_type is not None
            and content_type != expected_content_type
        ):
            raise InspectionError(
                "OOXML PowerPoint relationship type does not match its "
                "target content type"
            )
        expected_root = _PPT_RELATIONSHIP_TARGET_ROOTS.get(
            content_type or ""
        )
        if expected_root is None:
            if ooxml._is_xml_content_type(content_type):
                raise InspectionError(
                    "OOXML PowerPoint relationship target has an "
                    "unsupported XML content type"
                )
            continue
        if target in root_tags:
            continue
        root = ooxml._parse_graph_xml(
            archive,
            members,
            target,
            "PowerPoint relationship target",
        )
        if root.tag != expected_root:
            raise InspectionError(
                "OOXML PowerPoint relationship target root is invalid"
            )
        root_tags[target] = root.tag
    return root_tags


def _validate_word_relationship_targets(
    archive: zipfile.ZipFile,
    members: dict[str, zipfile.ZipInfo],
    part_content_types: dict[str, str],
    relationships: list[dict],
) -> None:
    for relationship in relationships:
        binding = _WORD_RELATIONSHIP_TARGET_ROOTS.get(relationship["type"])
        if binding is None:
            continue
        expected_content_type, expected_root = binding
        target = relationship.get("_normalized_target")
        content_type = (
            part_content_types.get(target)
            if isinstance(target, str)
            else None
        )
        if (
            relationship["external"]
            or not isinstance(target, str)
            or content_type != expected_content_type
        ):
            raise InspectionError(
                "OOXML Word relationship target content type is invalid"
            )
        root = ooxml._parse_graph_xml(
            archive,
            members,
            target,
            "Word relationship target",
        )
        if root.tag != expected_root:
            raise InspectionError(
                "OOXML Word relationship target root is invalid"
            )


def _validate_main_part_graph(
    archive: zipfile.ZipFile,
    members: dict[str, zipfile.ZipInfo],
    part_content_types: dict[str, str],
    relationships: list[dict],
    kind: str,
) -> None:
    main_part = _WORD_MAIN_PART if kind == "docx" else _PPT_MAIN_PART
    root = ooxml._parse_graph_xml(
        archive,
        members,
        main_part,
        "Office main",
    )
    if kind == "docx":
        document_tag = f"{{{_WORD_NAMESPACE}}}document"
        body_tag = f"{{{_WORD_NAMESPACE}}}body"
        if root.tag != document_tag:
            raise InspectionError("OOXML Word document root is invalid")
        bodies = [child for child in root if child.tag == body_tag]
        if len(bodies) != 1:
            raise InspectionError(
                "OOXML Word document body must appear exactly once"
            )
        _validate_word_relationship_targets(
            archive,
            members,
            part_content_types,
            relationships,
        )
        return

    presentation_tag = f"{{{_PRESENTATION_NAMESPACE}}}presentation"
    slide_tag = f"{{{_PRESENTATION_NAMESPACE}}}sld"
    notes_tag = f"{{{_PRESENTATION_NAMESPACE}}}notes"
    slide_list_tag = f"{{{_PRESENTATION_NAMESPACE}}}sldIdLst"
    slide_id_tag = f"{{{_PRESENTATION_NAMESPACE}}}sldId"
    relationship_id = f"{{{_OFFICE_RELATIONSHIPS_NAMESPACE}}}id"
    if root.tag != presentation_tag:
        raise InspectionError("OOXML PowerPoint presentation root is invalid")
    relationship_root_tags = _validate_ppt_relationship_target_roots(
        archive,
        members,
        part_content_types,
        relationships,
    )
    slide_lists = [child for child in root if child.tag == slide_list_tag]
    if len(slide_lists) > 1:
        raise InspectionError(
            "OOXML PowerPoint slide list must be unique"
        )
    slide_ids = list(slide_lists[0]) if slide_lists else []
    if any(element.tag != slide_id_tag for element in slide_ids):
        raise InspectionError("OOXML PowerPoint slide list is invalid")
    main_relationships = {
        row["id"]: row
        for row in relationships
        if row["source"] == main_part
    }
    main_slide_relationships = {
        row["id"]: row
        for row in relationships
        if (
            row["source"] == main_part
            and row["type"] == _SLIDE_RELATIONSHIP
        )
    }
    targets: set[str] = set()
    ids: set[str] = set()
    for slide_id in slide_ids:
        rel_id = slide_id.attrib.get(relationship_id)
        if not rel_id or rel_id in ids:
            raise InspectionError(
                "OOXML PowerPoint slide relationship id is invalid"
            )
        ids.add(rel_id)
        relationship = main_relationships.get(rel_id)
        target = (
            None
            if relationship is None
            else relationship.get("_normalized_target")
        )
        if (
            relationship is None
            or relationship["type"] != _SLIDE_RELATIONSHIP
            or relationship["external"]
            or not isinstance(target, str)
        ):
            raise InspectionError(
                "OOXML PowerPoint slide relationship is invalid"
            )
        if (
            target in targets
            or part_content_types.get(target) != _SLIDE_CONTENT_TYPE
        ):
            raise InspectionError(
                "OOXML PowerPoint slide relationship content type is invalid"
            )
        targets.add(target)

    slide_parts = {
        part
        for part, content_type in part_content_types.items()
        if content_type == _SLIDE_CONTENT_TYPE
    }
    if ids != set(main_slide_relationships) or targets != slide_parts:
        raise InspectionError(
            "OOXML PowerPoint slide relationship graph is incomplete"
        )

    notes_parts = {
        part
        for part, content_type in part_content_types.items()
        if content_type == _NOTES_SLIDE_CONTENT_TYPE
    }
    linked_notes: set[str] = set()
    typed_relationships = {
        _SLIDE_LAYOUT_RELATIONSHIP: (
            "slide layout",
            lambda content_type: content_type
            == _SLIDE_LAYOUT_CONTENT_TYPE,
        ),
        _NOTES_SLIDE_RELATIONSHIP: (
            "notes slide",
            lambda content_type: content_type
            == _NOTES_SLIDE_CONTENT_TYPE,
        ),
        _IMAGE_RELATIONSHIP: (
            "image",
            lambda content_type: content_type.startswith("image/"),
        ),
        _AUDIO_RELATIONSHIP: (
            "audio",
            lambda content_type: content_type.startswith("audio/"),
        ),
        _VIDEO_RELATIONSHIP: (
            "video",
            lambda content_type: content_type.startswith("video/"),
        ),
        _MEDIA_RELATIONSHIP: (
            "media",
            lambda content_type: content_type.startswith(
                ("image/", "audio/", "video/")
            ),
        ),
        _CHART_RELATIONSHIP: (
            "chart",
            lambda content_type: content_type in _CHART_CONTENT_TYPES,
        ),
    }
    for slide_part in sorted(slide_parts):
        if relationship_root_tags.get(slide_part) != slide_tag:
            raise InspectionError(
                "OOXML PowerPoint slide root is invalid"
            )
        slide_relationships = [
            row for row in relationships if row["source"] == slide_part
        ]
        layout_count = 0
        for relationship in slide_relationships:
            relationship_type = relationship["type"]
            validation = typed_relationships.get(relationship_type)
            if validation is None:
                continue
            label, accepts_content_type = validation
            target = relationship.get("_normalized_target")
            content_type = (
                part_content_types.get(target, "")
                if isinstance(target, str)
                else ""
            )
            if (
                relationship["external"]
                or not isinstance(target, str)
                or not accepts_content_type(content_type)
            ):
                raise InspectionError(
                    "OOXML PowerPoint slide relationship "
                    f"for {label} is invalid"
                )
            if relationship_type == _SLIDE_LAYOUT_RELATIONSHIP:
                layout_count += 1
            elif relationship_type == _NOTES_SLIDE_RELATIONSHIP:
                if target in linked_notes:
                    raise InspectionError(
                        "OOXML PowerPoint notes slide relationship is invalid"
                    )
                linked_notes.add(target)
        if layout_count != 1:
            raise InspectionError(
                "OOXML PowerPoint slide relationship to layout is invalid"
            )

    if linked_notes != notes_parts:
        raise InspectionError(
            "OOXML PowerPoint notes slide relationship graph is incomplete"
        )
    for notes_part in sorted(notes_parts):
        if relationship_root_tags.get(notes_part) != notes_tag:
            raise InspectionError(
                "OOXML PowerPoint notes slide root is invalid"
            )


def preflight_office(path: str, expected_kind: str | None = None) -> _PackageInfo:
    """Validate a DOCX/PPTX ZIP before python-docx/python-pptx opens it."""
    if expected_kind not in {None, "docx", "pptx"}:
        raise InspectionError("unsupported Office OOXML kind")
    descriptor = ooxml._open_stable_descriptor(path, "OOXML Office package")
    try:
        if not artifact_security._preflight_zip_descriptor(
            descriptor,
            validate_local_bounds=False,
        ):
            raise InspectionError("OOXML ZIP central directory failed preflight")
        with os.fdopen(os.dup(descriptor), "rb") as source:
            try:
                with zipfile.ZipFile(source) as archive:
                    members = archive.infolist()
                    if len(members) > ooxml._MAX_OOXML_ENTRIES:
                        raise InspectionError(
                            "OOXML entry count exceeds maximum "
                            f"{ooxml._MAX_OOXML_ENTRIES}"
                        )
                    names: set[str] = set()
                    members_by_name: dict[str, zipfile.ZipInfo] = {}
                    collision_keys: set[str] = set()
                    declared_total = 0
                    content_types_info = None
                    for member in members:
                        name = ooxml._validate_member_name(member)
                        collision_key = unicodedata.normalize(
                            "NFC", name
                        ).casefold()
                        if name in names or collision_key in collision_keys:
                            raise InspectionError(
                                "duplicate or colliding OOXML entry name"
                            )
                        names.add(name)
                        members_by_name[name] = member
                        collision_keys.add(collision_key)
                        declared_total += member.file_size
                        if declared_total > ooxml._MAX_OOXML_UNCOMPRESSED_BYTES:
                            raise InspectionError(
                                "OOXML declared uncompressed bytes exceed maximum "
                                f"{ooxml._MAX_OOXML_UNCOMPRESSED_BYTES}"
                            )
                        ooxml._validate_member_type(member)
                        if name == "[Content_Types].xml":
                            content_types_info = member
                    ooxml._validate_local_headers(descriptor, archive, members)
                    if content_types_info is None:
                        raise InspectionError("OOXML content types are missing")
                    content_types_payload = ooxml._read_bounded_member(
                        archive,
                        content_types_info,
                        ooxml._MAX_CONTENT_TYPES_BYTES,
                        "OOXML content types",
                    )
                    if expected_kind is None:
                        if _WORD_MAIN_PART in names:
                            kind = "docx"
                        elif _PPT_MAIN_PART in names:
                            kind = "pptx"
                        else:
                            raise InspectionError(
                                "OOXML Office main part is missing"
                            )
                    else:
                        kind = expected_kind
                    main_part, main_content_type = (
                        (_WORD_MAIN_PART, _WORD_MAIN_CONTENT_TYPE)
                        if kind == "docx"
                        else (_PPT_MAIN_PART, _PPT_MAIN_CONTENT_TYPE)
                    )
                    part_content_types = ooxml._validate_content_types(
                        content_types_payload,
                        names,
                        required_main_part=main_part,
                        required_main_content_type=main_content_type,
                    )

                    relationship_rows: list[dict] = []
                    totals: dict[str, int] = {}
                    actual_total = 0
                    media: list[dict] = []
                    embedded: list[dict] = []
                    macros: list[dict] = []
                    for member in members:
                        name = ooxml._member_name(member)
                        if member.is_dir():
                            if member.file_size:
                                raise InspectionError(
                                    "OOXML directory has invalid size metadata"
                                )
                            continue
                        content_type = part_content_types.get(name)
                        if name.endswith(".rels"):
                            if content_type != _RELATIONSHIPS_CONTENT_TYPE:
                                raise InspectionError(
                                    "OOXML relationship part content type is invalid"
                                )
                            source_part = _relationship_source(name)
                            if source_part and source_part not in names:
                                raise InspectionError(
                                    "OOXML relationship source part is missing"
                                )
                            relationship_rows.extend(
                                _parse_relationship_part(
                                    archive,
                                    member,
                                    source_part,
                                    names,
                                    totals,
                                )
                            )
                        xml_validator = (
                            _OfficeXmlValidator(name, content_type, totals)
                            if (
                                name == "[Content_Types].xml"
                                or ooxml._is_xml_content_type(content_type)
                            )
                            else None
                        )
                        actual_entry = 0
                        with archive.open(member, "r") as entry:
                            while True:
                                chunk = entry.read(_ZIP_STREAM_CHUNK)
                                if not chunk:
                                    break
                                actual_entry += len(chunk)
                                actual_total += len(chunk)
                                if (
                                    actual_entry > member.file_size
                                    or actual_total
                                    > ooxml._MAX_OOXML_UNCOMPRESSED_BYTES
                                ):
                                    raise InspectionError(
                                        "OOXML actual uncompressed bytes exceed "
                                        "declared or configured limits"
                                    )
                                if xml_validator is not None:
                                    xml_validator.feed(chunk)
                        if actual_entry != member.file_size:
                            raise InspectionError(
                                "OOXML entry size metadata does not match content"
                            )
                        if xml_validator is not None:
                            xml_validator.finish()
                        media_type = (content_type or "").casefold()
                        metadata = {
                            "part": _bounded_text(name, 1_024),
                            "content_type": _bounded_text(
                                content_type or "application/octet-stream",
                                512,
                            ),
                            "bytes": member.file_size,
                        }
                        if "vba" in media_type or "vba" in name.casefold():
                            if member.file_size > _MAX_MEDIA_PART_BYTES:
                                raise InspectionError(
                                    "OOXML macro part exceeds configured limit"
                                )
                            macros.append(metadata)
                        elif media_type.startswith(("image/", "audio/", "video/")):
                            if member.file_size > _MAX_MEDIA_PART_BYTES:
                                raise InspectionError(
                                    "OOXML media part exceeds configured limit"
                                )
                            media.append(metadata)
                        elif (
                            "oleobject" in media_type
                            or media_type.endswith(".package")
                            or "/embeddings/" in f"/{name.casefold()}"
                        ):
                            if member.file_size > _MAX_MEDIA_PART_BYTES:
                                raise InspectionError(
                                    "OOXML embedded object exceeds configured limit"
                                )
                            embedded.append(metadata)
                    if actual_total != declared_total:
                        raise InspectionError(
                            "OOXML total size metadata does not match content"
                        )

                    _validate_main_part_graph(
                        archive,
                        members_by_name,
                        part_content_types,
                        relationship_rows,
                        kind,
                    )
                    root_office = [
                        row
                        for row in relationship_rows
                        if row["source"] == "/"
                        and row["type"] == _OFFICE_DOCUMENT_RELATIONSHIP
                    ]
                    if (
                        len(root_office) != 1
                        or root_office[0]["external"]
                        or root_office[0]["_normalized_target"] != main_part
                    ):
                        raise InspectionError(
                            "OOXML package officeDocument relationship is invalid"
                        )
                    clean_relationships = tuple(
                        {
                            key: value
                            for key, value in row.items()
                            if not key.startswith("_")
                        }
                        for row in relationship_rows
                    )
                    return _PackageInfo(
                        kind=kind,
                        relationships=clean_relationships,
                        media=tuple(media),
                        embedded_objects=tuple(embedded),
                        macros=tuple(macros),
                    )
            except InspectionError:
                raise
            except (
                OSError,
                RuntimeError,
                NotImplementedError,
                ValueError,
                zlib.error,
                zipfile.BadZipFile,
                zipfile.LargeZipFile,
            ) as exc:
                raise InspectionError(
                    "OOXML ZIP could not be validated: "
                    f"{bounded_diagnostic(exc)}"
                ) from exc
    finally:
        os.close(descriptor)


def _is_ole(path: str) -> bool:
    descriptor = ooxml._open_stable_descriptor(path, "Office artifact")
    try:
        return os.pread(descriptor, len(_OLE_SIGNATURE), 0) == _OLE_SIGNATURE
    finally:
        os.close(descriptor)


def _ole_directory_names(path: str) -> set[str]:
    descriptor = ooxml._open_stable_descriptor(path, "legacy Office artifact")
    try:
        header = os.pread(descriptor, 512, 0)
        if len(header) != 512 or header[:8] != _OLE_SIGNATURE:
            raise InspectionError("legacy Office input has an invalid OLE header")
        sector_shift = struct.unpack_from("<H", header, 30)[0]
        sector_size = 1 << sector_shift
        if sector_size not in {512, 4096}:
            raise InspectionError("legacy Office input has an invalid sector size")
        file_size = os.fstat(descriptor).st_size
        sector_count = max(0, (file_size - 512) // sector_size)

        def read_sector(sector: int) -> bytes:
            if sector < 0 or sector >= sector_count:
                raise InspectionError(
                    "legacy Office input has an invalid sector chain"
                )
            payload = os.pread(
                descriptor,
                sector_size,
                512 + sector * sector_size,
            )
            if len(payload) != sector_size:
                raise InspectionError(
                    "legacy Office input has a truncated sector"
                )
            return payload

        fat_sectors = [
            value
            for value in struct.unpack_from("<109L", header, 76)
            if value not in {
                _OLE_FREESECT,
                _OLE_ENDOFCHAIN,
                _OLE_FATSECT,
                _OLE_DIFSECT,
            }
        ]
        fat: list[int] = []
        for sector in fat_sectors:
            fat.extend(
                struct.unpack(
                    f"<{sector_size // 4}L",
                    read_sector(sector),
                )
            )
        directory_sector = struct.unpack_from("<L", header, 48)[0]
        payload = bytearray()
        seen: set[int] = set()
        while directory_sector != _OLE_ENDOFCHAIN:
            if (
                directory_sector in seen
                or directory_sector >= len(fat)
                or len(seen) >= 128
            ):
                raise InspectionError(
                    "legacy Office input has an invalid directory chain"
                )
            seen.add(directory_sector)
            payload.extend(read_sector(directory_sector))
            directory_sector = fat[directory_sector]
        names: set[str] = set()
        for offset in range(0, len(payload), 128):
            entry = payload[offset:offset + 128]
            if len(entry) < 128:
                break
            name_bytes = struct.unpack_from("<H", entry, 64)[0]
            if name_bytes < 2 or name_bytes > 64 or name_bytes % 2:
                continue
            try:
                name = entry[:name_bytes - 2].decode("utf-16le")
            except UnicodeDecodeError:
                continue
            if name:
                names.add(name)
        return names
    finally:
        os.close(descriptor)


def _legacy_kind(path: str) -> str:
    names = _ole_directory_names(path)
    if "WordDocument" in names:
        return "doc"
    if "PowerPoint Document" in names:
        return "ppt"
    raise InspectionError("unsupported legacy Office compound document")


def _write_libreoffice_profile(profile: Path) -> None:
    user = profile / "user"
    user.mkdir(parents=True)
    config = """<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry">
  <item oor:path="/org.openoffice.Office.Common/Security/Scripting">
    <prop oor:name="MacroSecurityLevel" oor:op="fuse"><value>3</value></prop>
    <prop oor:name="DisableMacrosExecution" oor:op="fuse"><value>true</value></prop>
    <prop oor:name="DisableActiveContent" oor:op="fuse"><value>true</value></prop>
    <prop oor:name="BlockUntrustedRefererLinks" oor:op="fuse"><value>true</value></prop>
  </item>
  <item oor:path="/org.openoffice.Office.Common/Load">
    <prop oor:name="UpdateDocMode" oor:op="fuse"><value>0</value></prop>
  </item>
  <item oor:path="/org.openoffice.Office.Writer/Content/Update">
    <prop oor:name="Link" oor:op="fuse"><value>false</value></prop>
    <prop oor:name="Field" oor:op="fuse"><value>false</value></prop>
    <prop oor:name="Chart" oor:op="fuse"><value>false</value></prop>
  </item>
</oor:items>
"""
    (user / "registrymodifications.xcu").write_text(config, encoding="utf-8")
    fontconfig = """<?xml version="1.0"?>
<!DOCTYPE fontconfig SYSTEM "urn:fontconfig:fonts.dtd">
<fontconfig>
  <dir>/usr/share/fonts</dir>
  <dir>/usr/local/share/fonts</dir>
  <cachedir>/var/cache/fontconfig</cachedir>
  <include ignore_missing="yes">/etc/fonts/conf.d</include>
  <config><rescan><int>0</int></rescan></config>
</fontconfig>
"""
    (profile / "fontconfig.xml").write_text(fontconfig, encoding="utf-8")


def _find_libreoffice() -> str:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if executable is None:
        raise InspectionError("LibreOffice is not installed or not on PATH")
    return executable


def _convert_with_libreoffice(
    path: str,
    scratch_dir: str,
    target_format: str,
) -> str:
    if current_scratch_transaction() is None:
        raise InspectionError(
            "LibreOffice conversion requires an active scratch transaction"
        )
    if target_format not in {"docx", "pptx", "pdf"}:
        raise InspectionError("unsupported LibreOffice conversion target")
    libreoffice = _find_libreoffice()
    env_executable = shutil.which("env")
    if env_executable is None:
        raise InspectionError("env executable is unavailable for LibreOffice")
    profile = Path(tempfile.mkdtemp(prefix="lo-profile-", dir=scratch_dir))
    output_dir = Path(tempfile.mkdtemp(prefix="lo-output-", dir=scratch_dir))
    _write_libreoffice_profile(profile)
    descriptor = ooxml._open_stable_descriptor(
        path,
        "LibreOffice Office input",
    )
    input_path = f"/proc/self/fd/{descriptor}"
    command = [
        env_executable,
        f"FONTCONFIG_FILE={profile / 'fontconfig.xml'}",
        libreoffice,
        "--headless",
        "--invisible",
        "--nologo",
        "--nodefault",
        "--nolockcheck",
        "--norestore",
        f"-env:UserInstallation={profile.as_uri()}",
        "--convert-to",
        target_format,
        "--outdir",
        str(output_dir),
        input_path,
    ]
    try:
        deadline = time.monotonic() + _LIBREOFFICE_TIMEOUT_SECONDS
        for attempt in range(_LIBREOFFICE_MAX_TRANSIENT_ATTEMPTS):
            remaining_timeout = deadline - time.monotonic()
            if remaining_timeout <= 0:
                raise InspectionError(
                    "LibreOffice conversion exceeded its total timeout"
                )
            try:
                safe_run(
                    command,
                    timeout=remaining_timeout,
                    cwd=scratch_dir,
                    home=scratch_dir,
                    pass_fds=(descriptor,),
                )
                break
            except InspectionError as exc:
                if (
                    attempt + 1
                    >= _LIBREOFFICE_MAX_TRANSIENT_ATTEMPTS
                    or not ooxml._is_transient_libreoffice_scratch_error(
                        exc,
                        profile,
                    )
                ):
                    raise
        outputs = list(output_dir.iterdir())
        expected = output_dir / f"{descriptor}.{target_format}"
        if outputs != [expected]:
            raise InspectionError(
                "LibreOffice did not create one uniquely named output"
            )
        info = expected.lstat()
        if not stat.S_ISREG(info.st_mode) or info.st_nlink != 1:
            raise InspectionError(
                "LibreOffice output is not a single-link regular file"
            )
        if info.st_size <= 0:
            raise InspectionError("LibreOffice output is empty")
        with expected.open("rb") as source:
            signature = source.read(8)
        if target_format == "pdf":
            if not signature.startswith(b"%PDF-"):
                raise InspectionError("LibreOffice output is not a PDF")
        elif not signature.startswith(b"PK"):
            raise InspectionError("LibreOffice output is not an OOXML ZIP")
        if target_format in {"docx", "pptx"}:
            converted = preflight_office(str(expected), target_format)
            if converted.kind != target_format:
                raise InspectionError(
                    "LibreOffice output content type is incorrect"
                )
        return str(expected)
    except OSError as exc:
        raise InspectionError(
            f"LibreOffice output could not be validated: "
            f"{bounded_diagnostic(exc)}"
        ) from exc
    finally:
        os.close(descriptor)


def _prepare_office(path: str, scratch_dir: str) -> tuple[str, _PackageInfo, str | None]:
    if _is_ole(path):
        legacy_kind = _legacy_kind(path)
        target = "docx" if legacy_kind == "doc" else "pptx"
        converted = _convert_with_libreoffice(path, scratch_dir, target)
        return converted, preflight_office(converted, target), f".{legacy_kind}"
    info = preflight_office(path)
    return path, info, None


def _open_document(path: str):
    descriptor = ooxml._open_stable_descriptor(path, "Word document")
    try:
        with os.fdopen(os.dup(descriptor), "rb") as source:
            return Document(source)
    except Exception as exc:
        raise InspectionError(
            f"Word document could not be parsed: {bounded_diagnostic(exc)}"
        ) from exc
    finally:
        os.close(descriptor)


def _open_presentation(path: str):
    descriptor = ooxml._open_stable_descriptor(path, "PowerPoint presentation")
    try:
        with os.fdopen(os.dup(descriptor), "rb") as source:
            return Presentation(source)
    except Exception as exc:
        raise InspectionError(
            "PowerPoint presentation could not be parsed: "
            f"{bounded_diagnostic(exc)}"
        ) from exc
    finally:
        os.close(descriptor)


def _paragraph_detail(paragraph, index: int) -> dict:
    runs = []
    hyperlinks = []
    run_count = 0
    hyperlink_count = 0
    for item in paragraph.iter_inner_content():
        if hasattr(item, "address") and hasattr(item, "fragment"):
            hyperlink_count += 1
            if len(hyperlinks) < 64:
                hyperlinks.append(
                    {
                        "text": _bounded_text(item.text),
                        "address": _bounded_text(item.address, 1_024),
                        "fragment": _bounded_text(item.fragment, 512),
                    }
                )
            continue
        run_count += 1
        if len(runs) < 64:
            runs.append(
                {
                    "run": run_count,
                    "text": _bounded_text(item.text),
                    "bold": item.bold,
                    "italic": item.italic,
                    "underline": (
                        bool(item.underline)
                        if item.underline is not None
                        else None
                    ),
                }
            )
    return {
        "paragraph": index,
        "text": _bounded_text(paragraph.text),
        "style": _bounded_text(
            getattr(getattr(paragraph, "style", None), "name", ""), 256
        ),
        "runs": runs,
        "runs_omitted": max(0, run_count - len(runs)),
        "hyperlinks": hyperlinks,
        "hyperlinks_omitted": max(
            0, hyperlink_count - len(hyperlinks)
        ),
    }


def _nonempty_paragraphs(paragraphs) -> list[object]:
    return [
        paragraph
        for paragraph in paragraphs
        if str(getattr(paragraph, "text", "")).strip()
    ]


def _window_metadata(
    total: int,
    cursor: int,
    limit: int,
    returned: int,
) -> dict:
    next_cursor = cursor + returned
    if next_cursor >= total:
        next_cursor = None
    return {
        "cursor": cursor,
        "limit": limit,
        "total": total,
        "returned": returned,
        "omitted": max(0, total - cursor - returned),
        "next_cursor": next_cursor,
    }


def _validate_cursor(total: int, cursor: int, label: str) -> None:
    if cursor < 0 or cursor > total or (cursor == total and total > 0):
        raise InspectionError(f"{label} cursor exceeds item count")


def _cursor_items(
    items: list[object],
    cursor: int,
    limit: int,
    max_limit: int,
    build,
    budget: ResponseBudget,
    *,
    enforce_extract_chars: bool = False,
) -> dict:
    _validate_cursor(len(items), cursor, "Office selector")
    bounded_limit = min(limit, max_limit)
    result = {
        "items": [],
        **_window_metadata(len(items), cursor, bounded_limit, 0),
    }
    for index in range(cursor, min(len(items), cursor + bounded_limit)):
        item = build(index)
        returned = result["returned"] + 1
        candidate = {
            "items": [*result["items"], item],
            **_window_metadata(
                len(items),
                cursor,
                bounded_limit,
                returned,
            ),
        }
        if (
            _json_bytes(candidate) > max(0, budget.max_bytes // 2)
            or (
                enforce_extract_chars
                and _text_chars(candidate) > budget.max_extract_chars
            )
        ):
            break
        result = candidate
    return result


def _table_detail(
    table,
    index: int,
    *,
    selector_prefix: str | None = None,
    row_cursor: int = 0,
    row_limit: int = _TABLE_ROW_PAGE_SIZE,
    cell_row: int | None = None,
    cell_cursor: int = 0,
    cell_limit: int = _TABLE_CELL_PAGE_SIZE,
) -> dict:
    selector_prefix = selector_prefix or f"table:{index}"
    rows = list(table.rows)
    bounded_row_limit = min(row_limit, _TABLE_ROW_PAGE_SIZE)
    bounded_cell_limit = min(cell_limit, _TABLE_CELL_PAGE_SIZE)
    _validate_cursor(len(rows), row_cursor, "document table row")
    columns = len(table.columns)
    matrix: list[list[str]] = []
    cell_page = None
    if cell_row is not None:
        if cell_row > len(rows):
            raise InspectionError(
                "document table cell selector exceeds row count"
            )
        row_cursor = cell_row - 1
        row_cells = list(rows[row_cursor].cells)
        _validate_cursor(
            len(row_cells),
            cell_cursor,
            "document table cell",
        )
        selected_cells = row_cells[
            cell_cursor:cell_cursor + bounded_cell_limit
        ]
        matrix = [
            [_bounded_text(cell.text) for cell in selected_cells]
        ]
        cell_page = {
            "row": cell_row,
            **_window_metadata(
                len(row_cells),
                cell_cursor,
                bounded_cell_limit,
                len(selected_cells),
            ),
        }
        returned_rows = 1
        effective_row_limit = 1
    else:
        cell_count = 0
        returned_rows = 0
        for row in rows[
            row_cursor:row_cursor + bounded_row_limit
        ]:
            row_cells = list(row.cells)
            if cell_count + len(row_cells) > _TABLE_CELL_PAGE_SIZE:
                if returned_rows:
                    break
                selected_cells = row_cells[:_TABLE_CELL_PAGE_SIZE]
                matrix.append(
                    [
                        _bounded_text(cell.text)
                        for cell in selected_cells
                    ]
                )
                returned_rows = 1
                cell_page = {
                    "row": row_cursor + 1,
                    **_window_metadata(
                        len(row_cells),
                        0,
                        _TABLE_CELL_PAGE_SIZE,
                        len(selected_cells),
                    ),
                }
                break
            matrix.append(
                [_bounded_text(cell.text) for cell in row_cells]
            )
            cell_count += len(row_cells)
            returned_rows += 1
        effective_row_limit = bounded_row_limit
        if cell_page is None:
            returned_cells = sum(len(row) for row in matrix)
            cell_page = {
                "row": None,
                **_window_metadata(
                    returned_cells,
                    0,
                    _TABLE_CELL_PAGE_SIZE,
                    returned_cells,
                ),
            }
    row_page = _window_metadata(
        len(rows),
        row_cursor,
        effective_row_limit,
        returned_rows,
    )
    if row_page["next_cursor"] is not None:
        row_page["next_selector"] = (
            f"{selector_prefix}:rows:cursor:{row_page['next_cursor']}:"
            f"limit:{effective_row_limit}"
        )
    if cell_page["next_cursor"] is not None:
        cell_page["next_selector"] = (
            f"{selector_prefix}:row:{cell_page['row']}:cells:"
            f"cursor:{cell_page['next_cursor']}:"
            f"limit:{cell_page['limit']}"
        )
    return {
        "table": index,
        "rows": len(rows),
        "columns": columns,
        "cells": matrix,
        "row_page": row_page,
        "cell_page": cell_page,
    }


def _section_detail(section, index: int) -> dict:
    return {
        "section": index,
        "width": int(section.page_width or 0),
        "height": int(section.page_height or 0),
        "orientation": _bounded_text(section.orientation, 64),
        "margins": {
            "top": int(section.top_margin or 0),
            "right": int(section.right_margin or 0),
            "bottom": int(section.bottom_margin or 0),
            "left": int(section.left_margin or 0),
        },
    }


def _header_footer_parts(document, header: bool) -> list[object]:
    attributes = (
        ("header", "first_page_header", "even_page_header")
        if header
        else ("footer", "first_page_footer", "even_page_footer")
    )
    parts = []
    seen: set[str] = set()
    for section in document.sections:
        for attribute in attributes:
            item = getattr(section, attribute)
            paragraphs = _nonempty_paragraphs(item.paragraphs)
            tables = list(item.tables)
            if not paragraphs and not tables:
                continue
            key = str(item.part.partname)
            if key in seen:
                continue
            seen.add(key)
            parts.append(item)
    return parts


def _header_footer_detail(
    item,
    index: int,
    page: int,
    label: str,
    *,
    table_cursor: int = 0,
    table_limit: int = _TABLE_PAGE_SIZE,
    table_detail_selection: dict | None = None,
) -> dict:
    paragraphs = _nonempty_paragraphs(item.paragraphs)
    start = (page - 1) * _PARAGRAPH_PAGE_SIZE
    if start >= len(paragraphs) and (paragraphs or page > 1):
        raise InspectionError(
            f"document {label} page exceeds paragraph count"
        )
    selected = paragraphs[start:start + _PARAGRAPH_PAGE_SIZE]
    tables = list(item.tables)
    _validate_cursor(len(tables), table_cursor, f"document {label} table")
    bounded_table_limit = min(table_limit, _TABLE_PAGE_SIZE)
    selected_tables = tables[
        table_cursor:table_cursor + bounded_table_limit
    ]
    table_page = _window_metadata(
        len(tables),
        table_cursor,
        bounded_table_limit,
        len(selected_tables),
    )
    if table_page["next_cursor"] is not None:
        table_page["next_selector"] = (
            f"{label}:{index}:tables:"
            f"cursor:{table_page['next_cursor']}:"
            f"limit:{bounded_table_limit}"
        )
    return {
        label: index,
        "paragraphs": [_bounded_text(paragraph.text) for paragraph in selected],
        "tables": [
            _table_detail(
                table,
                table_cursor + offset,
                selector_prefix=(
                    f"{label}:{index}:table:{table_cursor + offset}"
                ),
                **(table_detail_selection or {}),
            )
            for offset, table in enumerate(
                selected_tables, start=1
            )
        ],
        "table_count": len(tables),
        "table_page": table_page,
        "paragraph_page": {
            "page": page,
            "total": len(paragraphs),
            "returned": len(selected),
            "omitted": len(paragraphs) - len(selected),
        },
    }


def _paged_items(
    total: int,
    page: int,
    page_size: int,
    build,
    budget: ResponseBudget,
    *,
    enforce_extract_chars: bool = False,
) -> dict:
    start = (page - 1) * page_size
    if start >= total and total:
        raise InspectionError("Office selector page exceeds item count")
    result = {
        "items": [],
        "page": page,
        "total": total,
        "returned": 0,
        "omitted": total,
    }
    for index in range(start, min(total, start + page_size)):
        item = build(index)
        candidate = {
            **result,
            "items": [*result["items"], item],
            "returned": result["returned"] + 1,
            "omitted": total - result["returned"] - 1,
        }
        if (
            _json_bytes(candidate) > max(0, budget.max_bytes // 2)
            or (
                enforce_extract_chars
                and _text_chars(candidate) > budget.max_extract_chars
            )
        ):
            break
        result = candidate
    return result


def _core_properties(document) -> dict:
    properties = document.core_properties
    names = (
        "title",
        "subject",
        "author",
        "keywords",
        "comments",
        "last_modified_by",
        "category",
        "content_status",
        "identifier",
        "language",
        "version",
        "created",
        "modified",
        "last_printed",
        "revision",
    )
    return {
        name: _json_scalar(getattr(properties, name, None))
        for name in names
        if getattr(properties, name, None) is not None
    }


def _document_structure(
    document,
    package: _PackageInfo,
    budget: ResponseBudget,
    converted_from: str | None,
    pages: dict[str, int],
    *,
    enforce_extract_chars: bool = False,
) -> dict:
    paragraphs = _nonempty_paragraphs(document.paragraphs)
    tables = list(document.tables)
    sections = list(document.sections)
    headers = _header_footer_parts(document, True)
    footers = _header_footer_parts(document, False)
    collections = {
        "paragraphs": (
            paragraphs,
            _PARAGRAPH_PAGE_SIZE,
            lambda index: _paragraph_detail(paragraphs[index], index + 1),
        ),
        "tables": (
            tables,
            _TABLE_PAGE_SIZE,
            lambda index: _table_detail(tables[index], index + 1),
        ),
        "headers": (
            headers,
            _HEADER_PAGE_SIZE,
            lambda index: _header_footer_detail(
                headers[index], index + 1, 1, "header"
            ),
        ),
        "footers": (
            footers,
            _HEADER_PAGE_SIZE,
            lambda index: _header_footer_detail(
                footers[index], index + 1, 1, "footer"
            ),
        ),
    }
    if enforce_extract_chars:
        detail_result = {
            "kind": "document",
            "format": "docx",
            "opens": True,
            "units_inspected": [],
        }
        for category, page in pages.items():
            items, page_size, build = collections[category]
            detail_result[category] = _paged_items(
                len(items),
                page,
                page_size,
                build,
                budget,
                enforce_extract_chars=True,
            )
            detail_result["units_inspected"].append(
                f"{category}:page:{page}"
            )
        if converted_from is not None:
            detail_result["converted_from"] = converted_from
        return validate_json_result(
            detail_result,
            budget,
            enforce_extract_chars=True,
        )

    result = {
        "kind": "document",
        "format": "docx",
        "opens": True,
        "section_count": len(sections),
        "sections": _cursor_items(
            sections,
            0,
            _MAX_INDEX_ITEMS,
            _MAX_INDEX_ITEMS,
            lambda index: {"section": index + 1},
            budget,
        ),
        "core_properties": _core_properties(document),
        "paragraphs": _paged_items(
            len(paragraphs),
            pages.get("paragraphs", 1),
            _PARAGRAPH_PAGE_SIZE,
            lambda index: _paragraph_detail(paragraphs[index], index + 1),
            budget,
            enforce_extract_chars=enforce_extract_chars,
        ),
        "tables": _paged_items(
            len(tables),
            pages.get("tables", 1),
            _TABLE_PAGE_SIZE,
            lambda index: _table_detail(tables[index], index + 1),
            budget,
            enforce_extract_chars=enforce_extract_chars,
        ),
        "headers": _paged_items(
            len(headers),
            pages.get("headers", 1),
            _HEADER_PAGE_SIZE,
            lambda index: _header_footer_detail(
                headers[index], index + 1, 1, "header"
            ),
            budget,
            enforce_extract_chars=enforce_extract_chars,
        ),
        "footers": _paged_items(
            len(footers),
            pages.get("footers", 1),
            _HEADER_PAGE_SIZE,
            lambda index: _header_footer_detail(
                footers[index], index + 1, 1, "footer"
            ),
            budget,
            enforce_extract_chars=enforce_extract_chars,
        ),
        "relationships": _metadata_page(list(package.relationships), budget),
        "media": _metadata_page(list(package.media), budget),
        "embedded_objects": _metadata_page(
            list(package.embedded_objects), budget
        ),
        "macros": _metadata_page(list(package.macros), budget),
    }
    if converted_from is not None:
        result["converted_from"] = converted_from
    return validate_json_result(
        result,
        budget,
        enforce_extract_chars=enforce_extract_chars,
    )


def _shape_type_name(shape) -> str:
    shape_type = getattr(shape, "shape_type", None)
    return _bounded_text(
        getattr(shape_type, "name", str(shape_type)),
        128,
    )


def _shape_text(shape) -> list[str]:
    if getattr(shape, "has_text_frame", False):
        return [
            _bounded_text(paragraph.text)
            for paragraph in shape.text_frame.paragraphs[
                :_MAX_SHAPE_TEXT_ITEMS
            ]
            if paragraph.text.strip()
        ]
    if getattr(shape, "has_table", False):
        text = []
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text.append(_bounded_text(cell.text))
                if len(text) >= _MAX_SHAPE_TEXT_ITEMS:
                    return text
        return text
    return []


def _shape_detail(
    shape,
    index: int,
    *,
    slide_index: int,
    path: tuple[int, ...],
    children_cursor: int = 0,
    children_limit: int = _SHAPE_PAGE_SIZE,
    expand_children: bool = True,
) -> dict:
    result = {
        "shape": index,
        "type": _shape_type_name(shape),
        "name": _bounded_text(getattr(shape, "name", ""), 256),
        "bounds": {
            "left": int(getattr(shape, "left", 0)),
            "top": int(getattr(shape, "top", 0)),
            "width": int(getattr(shape, "width", 0)),
            "height": int(getattr(shape, "height", 0)),
        },
        "text": _shape_text(shape),
    }
    if getattr(shape, "has_chart", False):
        chart = shape.chart
        result["chart"] = {
            "type": _bounded_text(
                getattr(chart.chart_type, "name", chart.chart_type), 128
            ),
            "series_count": len(chart.series),
        }
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        children = list(shape.shapes)
        bounded_limit = min(children_limit, _SHAPE_PAGE_SIZE)
        if expand_children:
            _validate_cursor(
                len(children),
                children_cursor,
                "presentation group child",
            )
            selected = children[
                children_cursor:children_cursor + bounded_limit
            ]
        else:
            selected = []
        nested = [
            _shape_detail(
                child,
                children_cursor + child_offset,
                slide_index=slide_index,
                path=(*path, children_cursor + child_offset),
                expand_children=False,
            )
            for child_offset, child in enumerate(selected, start=1)
        ]
        returned = len(nested)
        child_page = _window_metadata(
            len(children),
            children_cursor,
            bounded_limit,
            returned,
        )
        if not expand_children and children:
            child_page.update(
                {
                    "cursor": 0,
                    "returned": 0,
                    "omitted": len(children),
                    "next_cursor": 0,
                }
            )
        next_cursor = child_page["next_cursor"]
        if next_cursor is not None:
            shape_path = ".".join(str(part) for part in path)
            child_page["next_selector"] = (
                f"slide:{slide_index}:shape:{shape_path}:children:"
                f"cursor:{next_cursor}:limit:{bounded_limit}"
            )
        result["children"] = nested
        result["children_page"] = child_page
        result["children_omitted"] = child_page["omitted"]
    return result


def _slide_text(slide) -> list[str]:
    text = []
    pending = list(reversed(list(slide.shapes)))
    while pending:
        shape = pending.pop()
        text.extend(_shape_text(shape))
        if len(text) >= _MAX_INDEX_ITEMS:
            break
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            pending.extend(reversed(list(shape.shapes)))
    return text[:_MAX_INDEX_ITEMS]


def _notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        return ""
    return "\n".join(
        paragraph.text
        for paragraph in frame.paragraphs
        if paragraph.text.strip()
    )


def _notes_page(
    slide,
    cursor: int = 0,
    limit: int = _MAX_TEXT_CHARS,
) -> tuple[str, dict]:
    text = _notes_text(slide)
    _validate_cursor(len(text), cursor, "presentation notes")
    bounded_limit = min(limit, _MAX_TEXT_CHARS)
    selected = text[cursor:cursor + bounded_limit]
    page = _window_metadata(
        len(text),
        cursor,
        bounded_limit,
        len(selected),
    )
    page["truncated"] = page["omitted"] > 0
    return selected, page


def _shape_at_path(slide, path: tuple[int, ...]):
    shapes = list(slide.shapes)
    shape = None
    for depth, number in enumerate(path):
        if number > len(shapes):
            raise InspectionError(
                "presentation shape selector exceeds shape count"
            )
        shape = shapes[number - 1]
        if depth + 1 < len(path):
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.GROUP:
                raise InspectionError(
                    "presentation shape selector traverses a non-group shape"
                )
            shapes = list(shape.shapes)
    return shape


def _slide_detail(
    slide,
    index: int,
    selection: dict | None = None,
) -> dict:
    selection = selection or {"mode": "shapes"}
    shapes = list(slide.shapes)
    mode = selection.get("mode", "shapes")
    if mode == "children":
        path = selection["path"]
        shape = _shape_at_path(slide, path)
        if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.GROUP:
            raise InspectionError(
                "presentation children selector requires a group shape"
            )
        details = [
            _shape_detail(
                shape,
                path[-1],
                slide_index=index,
                path=path,
                children_cursor=selection["children_cursor"],
                children_limit=selection["children_limit"],
            )
        ]
        shape_window = _window_metadata(
            len(shapes),
            path[0] - 1,
            1,
            1,
        )
    elif mode == "notes":
        details = []
        shape_window = _window_metadata(
            len(shapes),
            0,
            _SHAPE_PAGE_SIZE,
            0,
        )
    else:
        shape_cursor = selection.get("shape_cursor", 0)
        shape_limit = min(
            selection.get("shape_limit", _SHAPE_PAGE_SIZE),
            _SHAPE_PAGE_SIZE,
        )
        _validate_cursor(
            len(shapes),
            shape_cursor,
            "presentation shape",
        )
        selected = shapes[
            shape_cursor:shape_cursor + shape_limit
        ]
        details = [
            _shape_detail(
                shape,
                shape_cursor + offset,
                slide_index=index,
                path=(shape_cursor + offset,),
            )
            for offset, shape in enumerate(selected, start=1)
        ]
        shape_window = _window_metadata(
            len(shapes),
            shape_cursor,
            shape_limit,
            len(details),
        )
        if shape_window["next_cursor"] is not None:
            shape_window["next_selector"] = (
                f"slide:{index}:shapes:"
                f"cursor:{shape_window['next_cursor']}:"
                f"limit:{shape_limit}"
            )
        if "shape_page" in selection:
            shape_window["page"] = selection["shape_page"]
    notes_cursor = (
        selection.get("notes_cursor", 0)
        if mode == "notes"
        else 0
    )
    notes_limit = (
        selection.get("notes_limit", _MAX_TEXT_CHARS)
        if mode == "notes"
        else _MAX_TEXT_CHARS
    )
    notes_text, notes_window = _notes_page(
        slide,
        notes_cursor,
        notes_limit,
    )
    if notes_window["next_cursor"] is not None:
        notes_window["next_selector"] = (
            f"slide:{index}:notes:"
            f"cursor:{notes_window['next_cursor']}:"
            f"limit:{notes_window['limit']}"
        )
    return {
        "slide": index,
        "text": _slide_text(slide) if mode == "shapes" else [],
        "notes_text": notes_text,
        "notes_page": notes_window,
        "shapes": details,
        "shape_page": shape_window,
    }


def _presentation_structure(
    presentation,
    package: _PackageInfo,
    budget: ResponseBudget,
    converted_from: str | None,
    selections: dict[int, dict] | None,
    *,
    enforce_extract_chars: bool = False,
) -> dict:
    slides = list(presentation.slides)
    if selections is None:
        selected = list(range(min(len(slides), _MAX_INDEX_ITEMS)))
    else:
        selected = [index - 1 for index in selections]
    result = {
        "kind": "presentation",
        "format": "pptx",
        "opens": True,
        "slides": {
            "items": [],
            "returned": 0,
            "omitted": len(slides),
        },
    }
    if enforce_extract_chars:
        result["units_inspected"] = []
    else:
        result.update(
            {
                "slide_count": len(slides),
                "dimensions": {
                    "width": int(presentation.slide_width),
                    "height": int(presentation.slide_height),
                },
                "relationships": _metadata_page(
                    list(package.relationships), budget
                ),
                "media": _metadata_page(list(package.media), budget),
                "embedded_objects": _metadata_page(
                    list(package.embedded_objects), budget
                ),
                "macros": _metadata_page(list(package.macros), budget),
            }
        )
    for slide_index in selected:
        selection = (
            None if selections is None else selections[slide_index + 1]
        )
        item = _slide_detail(
            slides[slide_index],
            slide_index + 1,
            selection,
        )
        candidate = {
            **result["slides"],
            "items": [*result["slides"]["items"], item],
            "returned": result["slides"]["returned"] + 1,
            "omitted": len(slides) - result["slides"]["returned"] - 1,
        }
        draft = {**result, "slides": candidate}
        if enforce_extract_chars:
            draft["units_inspected"] = [
                *result["units_inspected"],
                selection["selector"],
            ]
        if (
            _json_bytes(draft) > budget.max_bytes
            or (
                enforce_extract_chars
                and _text_chars(draft) > budget.max_extract_chars
            )
        ):
            break
        result = draft
    if converted_from is not None:
        result["converted_from"] = converted_from
    return validate_json_result(
        result,
        budget,
        enforce_extract_chars=enforce_extract_chars,
    )


def _doc_extract_pages(
    selectors: list[str],
) -> tuple[dict[str, int], list[dict], dict | None]:
    pages: dict[str, int] = {}
    exact: list[dict] = []
    exact_keys: set[str] = set()
    section_window = None
    if not selectors:
        return {
            "paragraphs": 1,
            "tables": 1,
            "headers": 1,
            "footers": 1,
        }, exact, section_window
    for selector in selectors:
        if selector in exact_keys:
            raise InspectionError(
                f"document selector repeats item {selector!r}"
            )
        section_match = _DOC_SECTIONS_CURSOR_RE.fullmatch(selector)
        if section_match is not None:
            if section_window is not None:
                raise InspectionError(
                    "Office selector repeats category 'sections'"
                )
            section_window = {
                "cursor": int(section_match.group(1)),
                "limit": int(section_match.group(2)),
                "selector": selector,
            }
            continue
        page_match = _DOC_PAGE_SELECTOR_RE.fullmatch(selector)
        if page_match is not None:
            category = page_match.group(1)
            if category in pages:
                raise InspectionError(
                    f"Office selector repeats category {category!r}"
                )
            pages[category] = int(page_match.group(2))
            continue
        header_tables_match = _DOC_HEADER_TABLES_CURSOR_RE.fullmatch(
            selector
        )
        if header_tables_match is not None:
            exact.append(
                {
                    "label": header_tables_match.group(1),
                    "number": int(header_tables_match.group(2)),
                    "page": 1,
                    "table_cursor": int(header_tables_match.group(3)),
                    "table_limit": int(header_tables_match.group(4)),
                    "selector": selector,
                }
            )
            exact_keys.add(selector)
            continue
        scoped_rows_match = _DOC_SCOPED_TABLE_ROWS_CURSOR_RE.fullmatch(
            selector
        )
        if scoped_rows_match is not None:
            exact.append(
                {
                    "label": scoped_rows_match.group(1),
                    "number": int(scoped_rows_match.group(2)),
                    "page": 1,
                    "table_cursor": int(scoped_rows_match.group(3)) - 1,
                    "table_limit": 1,
                    "table_detail_selection": {
                        "row_cursor": int(scoped_rows_match.group(4)),
                        "row_limit": int(scoped_rows_match.group(5)),
                    },
                    "selector": selector,
                }
            )
            exact_keys.add(selector)
            continue
        scoped_cells_match = _DOC_SCOPED_TABLE_CELLS_CURSOR_RE.fullmatch(
            selector
        )
        if scoped_cells_match is not None:
            exact.append(
                {
                    "label": scoped_cells_match.group(1),
                    "number": int(scoped_cells_match.group(2)),
                    "page": 1,
                    "table_cursor": int(scoped_cells_match.group(3)) - 1,
                    "table_limit": 1,
                    "table_detail_selection": {
                        "cell_row": int(scoped_cells_match.group(4)),
                        "cell_cursor": int(scoped_cells_match.group(5)),
                        "cell_limit": int(scoped_cells_match.group(6)),
                    },
                    "selector": selector,
                }
            )
            exact_keys.add(selector)
            continue
        table_rows_match = _DOC_TABLE_ROWS_CURSOR_RE.fullmatch(selector)
        if table_rows_match is not None:
            exact.append(
                {
                    "label": "table",
                    "number": int(table_rows_match.group(1)),
                    "row_cursor": int(table_rows_match.group(2)),
                    "row_limit": int(table_rows_match.group(3)),
                    "selector": selector,
                }
            )
            exact_keys.add(selector)
            continue
        table_cells_match = _DOC_TABLE_CELLS_CURSOR_RE.fullmatch(selector)
        if table_cells_match is not None:
            exact.append(
                {
                    "label": "table",
                    "number": int(table_cells_match.group(1)),
                    "cell_row": int(table_cells_match.group(2)),
                    "cell_cursor": int(table_cells_match.group(3)),
                    "cell_limit": int(table_cells_match.group(4)),
                    "selector": selector,
                }
            )
            exact_keys.add(selector)
            continue
        item_match = _DOC_ITEM_SELECTOR_RE.fullmatch(selector)
        if item_match is None:
            raise InspectionError(
                "document selector must be '<category>:page:<n>' or "
                "'<item>:<n>[:page:<n>]' or a cursor/limit selector"
            )
        exact_keys.add(selector)
        label = item_match.group(1)
        page = int(item_match.group(3) or 1)
        if label == "section" and page != 1:
            raise InspectionError(
                "document section selector does not accept a page"
            )
        selection = {
            "label": label,
            "number": int(item_match.group(2)),
            "page": page,
            "selector": selector,
        }
        if label == "table":
            selection.update(
                {
                    "row_cursor": (page - 1) * _TABLE_ROW_PAGE_SIZE,
                    "row_limit": _TABLE_ROW_PAGE_SIZE,
                }
            )
        exact.append(selection)
    if pages and (exact or section_window is not None):
        raise InspectionError(
            "document category and item selectors cannot be mixed"
        )
    if exact and section_window is not None:
        raise InspectionError(
            "document section and item selectors cannot be mixed"
        )
    return pages, exact, section_window


class OfficeInspector:
    """Inspect DOC/DOCX/PPT/PPTX artifacts within registry budgets."""

    def inspect(
        self,
        path: str,
        scratch_dir: str,
        *,
        response_budget: ResponseBudget,
    ):
        prepared, package, converted_from = _prepare_office(path, scratch_dir)
        if package.kind == "docx":
            return _document_structure(
                _open_document(prepared),
                package,
                response_budget,
                converted_from,
                {},
            )
        return _presentation_structure(
            _open_presentation(prepared),
            package,
            response_budget,
            converted_from,
            None,
        )

    def extract(
        self,
        path: str,
        scratch_dir: str,
        selectors: list[str],
        *,
        response_budget: ResponseBudget,
    ):
        prepared, package, converted_from = _prepare_office(path, scratch_dir)
        if package.kind == "pptx":
            presentation = _open_presentation(prepared)
            slide_count = len(presentation.slides)
            selected: dict[int, dict] = {}
            if not selectors:
                selectors = ["slide:1"] if slide_count else []
            for selector in selectors:
                match = _SLIDE_SELECTOR_RE.fullmatch(selector)
                if match is not None:
                    slide = int(match.group(1))
                    page = int(match.group(2) or 1)
                    selection = {
                        "mode": "shapes",
                        "shape_cursor": (page - 1) * _SHAPE_PAGE_SIZE,
                        "shape_limit": _SHAPE_PAGE_SIZE,
                        "shape_page": page,
                        "selector": selector,
                    }
                else:
                    match = _SLIDE_SHAPES_CURSOR_RE.fullmatch(selector)
                    if match is not None:
                        slide = int(match.group(1))
                        selection = {
                            "mode": "shapes",
                            "shape_cursor": int(match.group(2)),
                            "shape_limit": int(match.group(3)),
                            "selector": selector,
                        }
                    else:
                        match = _SLIDE_CHILDREN_CURSOR_RE.fullmatch(
                            selector
                        )
                        if match is not None:
                            slide = int(match.group(1))
                            selection = {
                                "mode": "children",
                                "path": tuple(
                                    int(part)
                                    for part in match.group(2).split(".")
                                ),
                                "children_cursor": int(match.group(3)),
                                "children_limit": int(match.group(4)),
                                "selector": selector,
                            }
                        else:
                            match = _SLIDE_NOTES_CURSOR_RE.fullmatch(
                                selector
                            )
                            if match is not None:
                                slide = int(match.group(1))
                                selection = {
                                    "mode": "notes",
                                    "notes_cursor": int(match.group(2)),
                                    "notes_limit": int(match.group(3)),
                                    "selector": selector,
                                }
                            else:
                                slide = 0
                                selection = {}
                if not selection:
                    raise InspectionError(
                        "presentation selector must be "
                        "'slide:<n>[:page:<n>]' or a supported "
                        "cursor/limit selector"
                    )
                if slide > slide_count:
                    raise InspectionError(
                        "presentation selector exceeds slide count"
                    )
                if slide in selected:
                    raise InspectionError(
                        f"presentation selector repeats slide {slide}"
                    )
                selected[slide] = selection
            return _presentation_structure(
                presentation,
                package,
                response_budget,
                converted_from,
                selected,
                enforce_extract_chars=True,
            )

        document = _open_document(prepared)
        pages, exact, section_window = _doc_extract_pages(selectors)
        if section_window is not None:
            sections = list(document.sections)
            result = {
                "kind": "document",
                "format": "docx",
                "opens": True,
                "sections": _cursor_items(
                    sections,
                    section_window["cursor"],
                    section_window["limit"],
                    _MAX_INDEX_ITEMS,
                    lambda index: _section_detail(
                        sections[index], index + 1
                    ),
                    response_budget,
                    enforce_extract_chars=True,
                ),
                "units_inspected": [section_window["selector"]],
            }
            if converted_from is not None:
                result["converted_from"] = converted_from
            return validate_json_result(
                result,
                response_budget,
                enforce_extract_chars=True,
            )
        if exact:
            paragraphs = _nonempty_paragraphs(document.paragraphs)
            tables = list(document.tables)
            headers = _header_footer_parts(document, True)
            footers = _header_footer_parts(document, False)
            sections = list(document.sections)
            collections = {
                "paragraph": paragraphs,
                "table": tables,
                "header": headers,
                "footer": footers,
                "section": sections,
            }
            result = {
                "kind": "document",
                "format": "docx",
                "opens": True,
                "details": [],
                "units_inspected": [],
                "omitted": 0,
            }
            for selection in exact:
                label = selection["label"]
                number = selection["number"]
                page = selection.get("page", 1)
                collection = collections[label]
                if number > len(collection):
                    raise InspectionError(
                        f"document selector exceeds {label} count"
                    )
                item = collection[number - 1]
                if label == "paragraph":
                    detail = _paragraph_detail(item, number)
                elif label == "table":
                    detail = _table_detail(
                        item,
                        number,
                        row_cursor=selection.get("row_cursor", 0),
                        row_limit=selection.get(
                            "row_limit", _TABLE_ROW_PAGE_SIZE
                        ),
                        cell_row=selection.get("cell_row"),
                        cell_cursor=selection.get("cell_cursor", 0),
                        cell_limit=selection.get(
                            "cell_limit", _TABLE_CELL_PAGE_SIZE
                        ),
                    )
                elif label == "section":
                    detail = _section_detail(item, number)
                else:
                    detail = _header_footer_detail(
                        item,
                        number,
                        page,
                        label,
                        table_cursor=selection.get("table_cursor", 0),
                        table_limit=selection.get(
                            "table_limit", _TABLE_PAGE_SIZE
                        ),
                        table_detail_selection=selection.get(
                            "table_detail_selection"
                        ),
                    )
                candidate = {
                    **result,
                    "details": [*result["details"], detail],
                    "units_inspected": [
                        *result["units_inspected"],
                        selection["selector"],
                    ],
                }
                if (
                    _json_bytes(candidate) > response_budget.max_bytes
                    or _text_chars(candidate)
                    > response_budget.max_extract_chars
                ):
                    result["omitted"] += 1
                    continue
                result = candidate
            if converted_from is not None:
                result["converted_from"] = converted_from
            return validate_json_result(
                result,
                response_budget,
                enforce_extract_chars=True,
            )
        return _document_structure(
            document,
            package,
            response_budget,
            converted_from,
            pages,
            enforce_extract_chars=True,
        )

    def render(
        self,
        path: str,
        scratch_dir: str,
        selectors: list[str],
        budget: RenderBudget,
    ) -> list[str]:
        for selector in selectors:
            if _PAGE_SELECTOR_RE.fullmatch(selector) is None:
                raise InspectionError(
                    "Office render selector must be 'page:<n>'"
                )
        if not _is_ole(path):
            preflight_office(path)
        pdf_path = _convert_with_libreoffice(path, scratch_dir, "pdf")
        try:
            from .pdf_image import PdfInspector
        except (ImportError, ModuleNotFoundError) as exc:
            raise InspectionError(
                "PDF renderer is unavailable for Office rendering"
            ) from exc
        return PdfInspector().render(pdf_path, scratch_dir, selectors, budget)


__all__ = ["OfficeInspector", "preflight_office"]
